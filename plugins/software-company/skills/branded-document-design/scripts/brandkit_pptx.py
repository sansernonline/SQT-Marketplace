"""
brandkit_pptx.py — สไลด์ PowerPoint โทนเดียวกับ brandkit.py (Apps Track style)
16:9 · พื้นขาว · แถบ accent ไล่สีน้ำเงิน-ม่วง · ตัวหนังสือ Tahoma

    pip install python-pptx

    from brandkit_pptx import BrandDeck

    d = BrandDeck()
    d.title_slide("Apps Track", "Project Control & Monitor",
                  "รายงานความคืบหน้า · 19 กรกฎาคม 2026")
    d.section("1 · ภาพรวมระบบ")
    d.bullets_slide("ขอบเขตงาน", ["โครงการและ Work items", "Sprint / Kanban board",
                                  "Timesheet และรายงาน"])
    d.kpi_slide("ตัวเลขสำคัญ", [("19", "โครงการ"), ("115", "Work items"),
                                ("103", "Open tasks"), ("14", "ผู้ใช้")])
    d.table_slide("สถานะ Milestone", ["รหัส", "รายการ", "สถานะ"],
                  [["M-01", "ออกแบบฐานข้อมูล", "เสร็จ"]])
    d.save("deck.pptx")
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from brandkit import TOKENS, FONT

W, H = Inches(13.333), Inches(7.5)          # 16:9
MARGIN = Inches(0.9)


def rgb(token):
    return RGBColor.from_string(TOKENS.get(token, token))


def _fmt(tf, size, color, bold=False, align=PP_ALIGN.LEFT, space_after=6):
    """ตั้งฟอนต์ทุกย่อหน้าใน text frame — ต้องตั้ง latin ผ่าน run.font.name
    และปล่อยให้ PowerPoint ใช้ฟอนต์เดียวกับข้อความไทย (Tahoma มีทั้งสองชุด)"""
    for p in tf.paragraphs:
        p.alignment = align
        p.space_after = Pt(space_after)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)


class BrandDeck:
    def __init__(self, template=None):
        self.prs = Presentation(template) if template else Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]

    # -- ชิ้นส่วนพื้นฐาน ---------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _rect(self, slide, x, y, w, h, fill=None, line=None):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill:
            shp.fill.solid()
            shp.fill.fore_color.rgb = rgb(fill)
        else:
            shp.fill.background()
        if line:
            shp.line.color.rgb = rgb(line)
            shp.line.width = Pt(0.75)
        else:
            shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _text(self, slide, text, x, y, w, h, size, color, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        lines = text.split("\n")
        tf.text = lines[0]
        for extra in lines[1:]:
            tf.add_paragraph().text = extra
        _fmt(tf, size, color, bold, align)
        return box

    def _accent_bar(self, slide, y=Emu(0), h=Inches(0.09)):
        """แถบ accent บาง ๆ ด้านบนสไลด์ — ตัวเชื่อมสายตาให้ทั้งเด็คดูเป็นชุดเดียว"""
        return self._rect(slide, Emu(0), y, W, h, fill="brand")

    def _footer(self, slide, left_text=""):
        n = len(self.prs.slides._sldIdLst)
        self._text(slide, left_text, MARGIN, H - Inches(0.62),
                   Inches(8), Inches(0.3), 10, "text_faint")
        self._text(slide, str(n), W - MARGIN - Inches(0.6), H - Inches(0.62),
                   Inches(0.6), Inches(0.3), 10, "text_faint", align=PP_ALIGN.RIGHT)

    # -- ประเภทสไลด์ -------------------------------------------------------
    def title_slide(self, title, subtitle=None, meta=None):
        s = self._slide()
        self._rect(s, Emu(0), Emu(0), W, H, fill="brand_tint")
        self._rect(s, Emu(0), Inches(3.05), W, Inches(0.06), fill="brand")
        self._text(s, title, MARGIN, Inches(2.05), W - 2 * MARGIN, Inches(1),
                   40, "brand_deep", bold=True)
        if subtitle:
            self._text(s, subtitle, MARGIN, Inches(3.3), W - 2 * MARGIN,
                       Inches(0.6), 20, "text", bold=True)
        if meta:
            self._text(s, meta, MARGIN, Inches(4.0), W - 2 * MARGIN,
                       Inches(0.5), 13, "text_muted")
        return s

    def section(self, title, kicker=None):
        s = self._slide()
        self._rect(s, Emu(0), Emu(0), Inches(0.28), H, fill="brand")
        if kicker:
            self._text(s, kicker, Inches(1.1), Inches(2.9), Inches(10),
                       Inches(0.4), 13, "brand_2", bold=True)
        self._text(s, title, Inches(1.1), Inches(3.3), Inches(11),
                   Inches(1), 32, "brand_deep", bold=True)
        return s

    def _content_head(self, title, subtitle=None):
        s = self._slide()
        self._accent_bar(s)
        self._text(s, title, MARGIN, Inches(0.62), W - 2 * MARGIN, Inches(0.6),
                   26, "brand", bold=True)
        y = Inches(1.35)
        if subtitle:
            self._text(s, subtitle, MARGIN, Inches(1.2), W - 2 * MARGIN,
                       Inches(0.4), 13, "text_muted")
            y = Inches(1.75)
        # ขีดสั้นใต้หัวข้อ — อ่านว่า "ตั้งใจ" กว่าเส้นคั่นเต็มความกว้าง
        self._rect(s, MARGIN, y - Inches(0.18), Inches(1.15), Inches(0.045),
                   fill="brand_2")
        return s, y

    def bullets_slide(self, title, items, subtitle=None):
        s, y = self._content_head(title, subtitle)
        box = s.shapes.add_textbox(MARGIN, y, W - 2 * MARGIN, H - y - Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "•  " + it
            p.space_after = Pt(12)
        _fmt(tf, 17, "text_body", space_after=12)
        self._footer(s)
        return s

    def kpi_slide(self, title, items, subtitle=None):
        s, y = self._content_head(title, subtitle)
        n = len(items)
        gap = Inches(0.25)
        total = W - 2 * MARGIN
        cw = int((total - gap * (n - 1)) / n)
        for i, (value, label) in enumerate(items):
            x = MARGIN + i * (cw + gap)
            self._rect(s, x, y + Inches(0.3), Emu(cw), Inches(1.9),
                       fill="brand_tint")
            self._text(s, str(value), x, y + Inches(0.6), Emu(cw), Inches(0.9),
                       40, "brand", bold=True, align=PP_ALIGN.CENTER)
            self._text(s, label, x, y + Inches(1.55), Emu(cw), Inches(0.5),
                       13, "text_muted", align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def table_slide(self, title, headers, rows, col_widths=None, subtitle=None,
                    status_col=None, palette=None):
        s, y = self._content_head(title, subtitle)
        total = W - 2 * MARGIN
        shape = s.shapes.add_table(len(rows) + 1, len(headers), MARGIN,
                                   y + Inches(0.25), total,
                                   Inches(0.42) * (len(rows) + 1))
        tbl = shape.table
        tbl.first_row = True
        if col_widths:
            unit = total / sum(col_widths)
            for i, w in enumerate(col_widths):
                tbl.columns[i].width = Emu(int(w * unit))
        for ci, h in enumerate(headers):
            c = tbl.cell(0, ci)
            c.text = str(h)
            c.fill.solid()
            c.fill.fore_color.rgb = rgb("brand_tint")
            c.margin_left = c.margin_right = Inches(0.12)
            _fmt(c.text_frame, 13, "brand_deep", bold=True, space_after=0)
        palette = palette or {}
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = "" if val is None else str(val)
                c.margin_left = c.margin_right = Inches(0.12)
                tone = palette.get(str(val).strip()) if ci == status_col else None
                c.fill.solid()
                c.fill.fore_color.rgb = rgb(tone + "_bg" if tone else "FFFFFF")
                _fmt(c.text_frame, 12.5, tone or "text_body", bold=bool(tone),
                     align=PP_ALIGN.CENTER if tone else PP_ALIGN.LEFT,
                     space_after=0)
        self._footer(s)
        return s

    def image_slide(self, title, image_path, caption=None, subtitle=None):
        s, y = self._content_head(title, subtitle)
        top = y + Inches(0.3)
        max_h = H - top - Inches(1.1)
        pic = s.shapes.add_picture(image_path, MARGIN, top, height=Emu(int(max_h)))
        if pic.width > W - 2 * MARGIN:      # กว้างเกิน → ยึดความกว้างแทน
            sp = s.shapes._spTree
            sp.remove(pic._element)
            pic = s.shapes.add_picture(image_path, MARGIN, top,
                                       width=Emu(int(W - 2 * MARGIN)))
        pic.left = Emu(int((W - pic.width) / 2))
        if caption:
            self._text(s, caption, MARGIN, H - Inches(1.0), W - 2 * MARGIN,
                       Inches(0.4), 12, "text_muted", align=PP_ALIGN.CENTER)
        self._footer(s)
        return s

    def quote_slide(self, text, source=None):
        s = self._slide()
        self._rect(s, Emu(0), Emu(0), W, H, fill="brand_tint")
        self._text(s, "“" + text + "”", Inches(1.6), Inches(2.4),
                   W - Inches(3.2), Inches(2.2), 26, "brand_deep", bold=True,
                   align=PP_ALIGN.CENTER)
        if source:
            self._text(s, "— " + source, Inches(1.6), Inches(4.6),
                       W - Inches(3.2), Inches(0.5), 14, "text_muted",
                       align=PP_ALIGN.CENTER)
        return s

    def save(self, path):
        self.prs.save(path)
        return path
